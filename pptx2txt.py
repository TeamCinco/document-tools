import os
from pptx import Presentation

def extract_text_from_pptx(file_path):
    """
    Extracts text from a PowerPoint file.
    """
    prs = Presentation(file_path)
    text_content = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

    return "\n".join(text_content)

def convert_pptx_to_txt(folder_path):
    """
    Converts all .pptx files in the specified folder to .txt files.
    """
    # List all files in the directory
    files = os.listdir(folder_path)
    
    # Process each .pptx file
    for file in files:
        if file.endswith(".pptx"):
            # Construct full file path
            file_path = os.path.join(folder_path, file)
            
            # Extract text from the pptx file
            text = extract_text_from_pptx(file_path)
            
            # Define output txt file name
            output_file_name = file.replace(".pptx", ".txt")
            output_file_path = os.path.join(folder_path, output_file_name)
            
            # Write extracted text to the txt file
            with open(output_file_path, "w", encoding="utf-8") as txt_file:
                txt_file.write(text)
    
    print("Conversion completed. All .pptx files have been converted to .txt files.")

# Usage: Specify the folder containing the .pptx files
folder_path = r"C:\Users\cinco\Desktop\pptx"
convert_pptx_to_txt(folder_path)
